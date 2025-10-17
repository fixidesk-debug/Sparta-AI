"""
Export Service - PDF, PowerPoint, Word Export
"""
from typing import Dict, Any, List, Optional
from io import BytesIO
import logging

logger = logging.getLogger(__name__)


class ExportService:
    """Export reports to multiple formats"""
    
    @staticmethod
    def export_to_pdf(
        title: str,
        content: str,
        charts: List[Dict[str, Any]] = None
    ) -> bytes:
        """Export to PDF using reportlab"""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import inch
            
            buffer = BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Title
            story.append(Paragraph(title, styles['Title']))
            story.append(Spacer(1, 0.2*inch))
            
            # Content
            for line in content.split('\n'):
                if line.strip():
                    story.append(Paragraph(line, styles['Normal']))
                    story.append(Spacer(1, 0.1*inch))
            
            doc.build(story)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"PDF export error: {e}")
            raise
    
    @staticmethod
    def export_to_powerpoint(
        title: str,
        slides_content: List[Dict[str, Any]]
    ) -> bytes:
        """Export to PowerPoint using python-pptx"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title
            
            # Content slides
            for slide_data in slides_content:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data.get('title', '')
                
                content = slide.placeholders[1].text_frame
                content.text = slide_data.get('content', '')
            
            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"PowerPoint export error: {e}")
            raise
    
    @staticmethod
    def export_to_word(
        title: str,
        content: str,
        tables: List[Dict[str, Any]] = None
    ) -> bytes:
        """Export to Word using python-docx"""
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            
            doc = Document()
            
            # Title
            doc.add_heading(title, 0)
            
            # Content
            for line in content.split('\n'):
                if line.strip().startswith('#'):
                    level = line.count('#')
                    doc.add_heading(line.replace('#', '').strip(), level)
                elif line.strip():
                    doc.add_paragraph(line)
            
            # Tables
            if tables:
                for table_data in tables:
                    doc.add_heading(table_data.get('title', 'Table'), 2)
                    rows = table_data.get('rows', [])
                    if rows:
                        table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                        for i, row in enumerate(rows):
                            for j, cell in enumerate(row):
                                table.rows[i].cells[j].text = str(cell)
            
            buffer = BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Word export error: {e}")
            raise
    
    @staticmethod
    def export_to_excel(
        df,
        filename: str,
        include_charts: bool = False
    ) -> bytes:
        """Export DataFrame to Excel with formatting"""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
            from openpyxl.utils.dataframe import dataframe_to_rows
            
            buffer = BytesIO()
            
            # Create workbook
            wb = Workbook()
            ws = wb.active
            ws.title = "Data"
            
            # Write headers with formatting
            header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF")
            
            for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True), 1):
                for c_idx, value in enumerate(row, 1):
                    cell = ws.cell(row=r_idx, column=c_idx, value=value)
                    
                    # Format header row
                    if r_idx == 1:
                        cell.fill = header_fill
                        cell.font = header_font
                        cell.alignment = Alignment(horizontal="center", vertical="center")
            
            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
            
            # Add summary sheet
            ws_summary = wb.create_sheet("Summary")
            ws_summary['A1'] = "Dataset Summary"
            ws_summary['A1'].font = Font(bold=True, size=14)
            
            ws_summary['A3'] = "Total Rows:"
            ws_summary['B3'] = len(df)
            ws_summary['A4'] = "Total Columns:"
            ws_summary['B4'] = len(df.columns)
            ws_summary['A5'] = "Filename:"
            ws_summary['B5'] = filename
            
            # Add statistics sheet
            ws_stats = wb.create_sheet("Statistics")
            ws_stats['A1'] = "Column Statistics"
            ws_stats['A1'].font = Font(bold=True, size=14)
            
            stats_df = df.describe(include='all').transpose()
            for r_idx, row in enumerate(dataframe_to_rows(stats_df, index=True, header=True), 3):
                for c_idx, value in enumerate(row, 1):
                    ws_stats.cell(row=r_idx, column=c_idx, value=value)
            
            wb.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Excel export error: {e}")
            raise
    
    @staticmethod
    def export_chart_to_png(
        chart_config: Dict[str, Any]
    ) -> bytes:
        """Export chart configuration to PNG image"""
        try:
            import matplotlib
            matplotlib.use('Agg')  # Use non-interactive backend
            import matplotlib.pyplot as plt
            import numpy as np
            
            # Extract chart configuration
            chart_type = chart_config.get('type', 'bar')
            data = chart_config.get('data', [])
            x_key = chart_config.get('xKey', 'x')
            y_key = chart_config.get('yKey', 'y')
            title = chart_config.get('title', 'Chart')
            
            # Create figure
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Extract data
            if data:
                x_values = [item.get(x_key, '') for item in data]
                y_values = [item.get(y_key, 0) for item in data]
                
                # Create chart based on type
                if chart_type == 'bar':
                    ax.bar(x_values, y_values, color='#3b82f6')
                elif chart_type == 'line':
                    ax.plot(x_values, y_values, marker='o', color='#3b82f6', linewidth=2)
                elif chart_type == 'scatter':
                    ax.scatter(x_values, y_values, color='#3b82f6', s=100, alpha=0.6)
                elif chart_type == 'pie':
                    ax.pie(y_values, labels=x_values, autopct='%1.1f%%', startangle=90)
                    ax.axis('equal')
                else:
                    ax.bar(x_values, y_values, color='#3b82f6')
                
                # Set labels and title
                if chart_type != 'pie':
                    ax.set_xlabel(x_key.capitalize())
                    ax.set_ylabel(y_key.capitalize())
                    ax.grid(True, alpha=0.3)
                
                ax.set_title(title, fontsize=14, fontweight='bold')
                
                # Rotate x-axis labels if needed
                if chart_type in ['bar', 'line'] and len(x_values) > 5:
                    plt.xticks(rotation=45, ha='right')
            
            plt.tight_layout()
            
            # Save to buffer
            buffer = BytesIO()
            plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
            plt.close(fig)
            
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"PNG export error: {e}")
            raise
